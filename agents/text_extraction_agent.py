
import os
import sys
import PyPDF2
from docx import Document
from pptx import Presentation

# Thêm đường dẫn cha để import module cấu hình model
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from langchain_core.tools import tool
from langgraph.prebuilt import create_react_agent
from langchain_core.messages import AIMessage
from langgraph.checkpoint.memory import MemorySaver
from langgraph.prebuilt import create_react_agent
from config.llm import gemini
from config.prompt import text_extraction_prompt
from agents.base import BaseAgent
from schemas.agent_schema import ResponseFormat
from utils.get_agent_response import get_agent_response
from typing import Any, AsyncIterable, Dict

# MemorySaver
memory = MemorySaver()


# ---------------------- TOOLS ----------------------
@tool("extract_text_from_pdf")
def extract_text_from_pdf(pdf_path: str) -> str:
    """Use this tool to extract text from a PDF file. Input is a string path to a .pdf file."""
    text = ""
    try:
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() or ""
        return text if text else "No extractable text found in PDF."
    except Exception as e:
        return f"Error extracting PDF: {e}"

@tool("extract_text_from_word")
def extract_text_from_word(word_path: str) -> str:
    """Use this tool to extract text from a Word (.docx) file. Input is a string path to a .docx file."""
    try:
        doc = Document(word_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text if text else "No text found in Word document."
    except Exception as e:
        return f"Error extracting Word: {e}"

@tool("extract_text_from_powerpoint")
def extract_text_from_powerpoint(ppt_path: str) -> str:
    """Use this tool to extract text from a PowerPoint (.pptx) file. Input is a string path to a .pptx file."""
    try:
        prs = Presentation(ppt_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text if text else "No text found in PowerPoint slides."
    except Exception as e:
        return f"Error extracting PowerPoint: {e}"

# ---------------------- AGENT SETUP ----------------------


class TextExtractionAgent(BaseAgent):
    """Text Extraction Agent backed by LangGraph."""

    def __init__(self):

        super().__init__(
            agent_name='TextExtractionAgent',
            description='Extract text from various file types such as PDF, Word, and PowerPoint',
            content_types=['text', 'text/plain'],
        )

        self.model = gemini

        self.graph = create_react_agent(
            self.model,
            checkpointer=memory,
            prompt=text_extraction_prompt,
            response_format=ResponseFormat,
            tools=[
                extract_text_from_pdf,
                extract_text_from_word,
                extract_text_from_powerpoint
            ],
            name="Text Extraction Agent",
        )

    def invoke(self, query, sessionId) -> str:
        config = {'configurable': {'thread_id': sessionId}, 'recursion_limit': 50}
        response = self.graph.invoke({'messages': [('user', query)]}, config)
        return get_agent_response(self.graph, response['messages'][-1], config)
    
    
    async def stream(
            self, query, sessionId, task_id
        ) -> AsyncIterable[Dict[str, Any]]:
            inputs = {'messages': [('user', query)]}
            config = {'configurable': {'thread_id': sessionId}, 'recursion_limit': 50}


            for item in self.graph.stream(inputs, config, stream_mode='values'):
                message = item['messages'][-1]
                if isinstance(message, AIMessage):
                    yield {
                        'response_type': 'text',
                        'is_task_complete': False,
                        'require_user_input': False,
                        'content': message.content,
                    }
            yield get_agent_response(self.graph, message, config)



# ---------------------- RUN TEST ----------------------


if __name__ == "__main__":
    agent = TextExtractionAgent()
    result = agent.invoke(query = "Extract text from the file at path C:/Users/dhuu3/Desktop/Chatbot_MCP/data/project-dv.docx", sessionId = "123")
    print(result)
