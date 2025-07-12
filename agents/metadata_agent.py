import os
import sys
import json
import asyncio
import mcp
import pandas as pd
from datetime import datetime  # Sửa lại import datetime
import subprocess
import time
import logging
import threading
from concurrent.futures import ThreadPoolExecutor
import queue
import requests
import re
from typing import Any, AsyncIterable, Dict, Optional, List, Union
import uuid



# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Đảm bảo đường dẫn đúng
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from langchain_core.tools import tool
from langgraph.checkpoint.memory import MemorySaver
from langgraph.prebuilt import create_react_agent
from langchain.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.messages import AIMessage
from config.llm import gemini
from agents.base import BaseAgent
from schemas.agent_schema import ResponseFormat
from utils import _format_file_timestamp
from config.prompt import metadata_prompt
from utils.get_agent_response import get_agent_response
from typing import Any, AsyncIterable, Dict, Optional, List

# MCP Client imports
try:
    from mcp import ClientSession
    from mcp.client.stdio import StdioServerParameters, stdio_client
    import mcp.types as types
    MCP_AVAILABLE = True
except ImportError:
    print("Warning: MCP client not available. Install with: pip install mcp")
    MCP_AVAILABLE = False

# MemorySaver
memory = MemorySaver()

class MCPConnection:
    """Manages MCP server connection with improved error handling and thread safety"""
    
    def __init__(self):
        self.server_params = None
        self.server_process = None
        self.is_connected = False
        self._lock = threading.Lock()
        self._connection_pool = queue.Queue(maxsize=5)
        self._pool_initialized = False
        
    def _run_async_in_thread(self, coro):
        """Run async function in a separate thread with its own event loop"""
        def run_in_new_loop():
            try:
                loop = asyncio.new_event_loop()
                asyncio.set_event_loop(loop)
                try:
                    return loop.run_until_complete(coro)
                finally:
                    loop.close()
            except Exception as e:
                logger.error(f"Error in async thread: {e}")
                return {"error": str(e)}
        
        with ThreadPoolExecutor(max_workers=1) as executor:
            future = executor.submit(run_in_new_loop)
            try:
                return future.result(timeout=30)
            except Exception as e:
                logger.error(f"Thread execution error: {e}")
                return {"error": str(e)}
    
    async def start_server(self):
        """Start MCP server with improved connection handling"""
        try:
            with self._lock:
                if self.is_connected:
                    return True
            
            # Get the absolute path to the MCP server
            current_dir = os.path.dirname(os.path.abspath(__file__))
            project_root = os.path.dirname(current_dir)
            server_path = os.path.join(project_root, 'fast_mcp', 'mcp_metadata_server.py')
            
            print(f"🔍 Looking for MCP server at: {server_path}")
            
            if not os.path.exists(server_path):
                # Try alternative paths
                alt_paths = [
                    os.path.join(os.path.dirname(__file__), 'mcp_metadata_server.py'),
                    os.path.join(project_root, 'mcp_metadata_server.py'),
                    os.path.join(current_dir, 'mcp_metadata_server.py')
                ]
                
                for alt_path in alt_paths:
                    if os.path.exists(alt_path):
                        server_path = alt_path
                        print(f"📍 Found server at: {server_path}")
                        break
                else:
                    raise FileNotFoundError(f"MCP server not found at any of: {[server_path] + alt_paths}")
            
            # Create server parameters
            params = StdioServerParameters(
                command="python",
                args=[server_path],
                env=os.environ.copy(),
                cwd=project_root
            )
            
            # Test connection
            try:
                print("🔄 Testing MCP server connection...")
                async with asyncio.timeout(15):
                    async with stdio_client(params) as streams:
                        async with mcp.ClientSession(*streams) as session:
                            await session.initialize()
                            
                            # List available tools to verify connection
                            tools = await session.list_tools()
                            tool_names = [tool.name for tool in tools.tools]
                            print(f"✅ Connected! Available tools: {tool_names}")
                            
                            # Test a simple tool call if available
                            if "list_all_metadata" in tool_names:
                                try:
                                    test_result = await session.call_tool("list_all_metadata", {})
                                    print(f"✅ Test call successful")
                                except Exception as test_e:
                                    print(f"⚠️ Test call failed but connection OK: {test_e}")
                
                with self._lock:
                    self.server_params = params
                    self.is_connected = True
                
                print("✅ MCP server connection established successfully")
                return True
                
            except asyncio.TimeoutError:
                print("❌ Connection timeout - server failed to start within 15 seconds")
                return False
            except Exception as e:
                print(f"❌ Connection test failed: {e}")
                import traceback
                print(traceback.format_exc())
                return False
            
        except Exception as e:
            print(f"❌ Failed to establish MCP connection: {e}")
            import traceback
            print(traceback.format_exc())
            return False
    
    def call_tool_sync(self, tool_name: str, arguments: dict):
        """Synchronous wrapper for MCP tool calls"""
        return self._run_async_in_thread(self._call_tool_async(tool_name, arguments))
    
    async def _call_tool_async(self, tool_name: str, arguments: dict):
        """Async implementation of tool call"""
        if not self.server_params:
            raise RuntimeError("MCP server not initialized. Call start_server() first.")
            
        max_retries = 2
        retry_delay = 1
        
        for attempt in range(max_retries):
            try:
                async with asyncio.timeout(30):
                    async with stdio_client(self.server_params) as streams:
                        async with mcp.ClientSession(*streams) as session:
                            await session.initialize()
                            
                            # Call the tool
                            result = await session.call_tool(tool_name, arguments)
                            
                            # Parse the result
                            if result and hasattr(result, 'content') and result.content:
                                content_item = result.content[0]
                                if hasattr(content_item, 'text'):
                                    content = content_item.text
                                else:
                                    content = str(content_item)
                                
                                try:
                                    # Try to parse as JSON
                                    return json.loads(content) if content else {"error": "Empty response"}
                                except json.JSONDecodeError:
                                    # Return as is if not valid JSON
                                    return {"result": content}
                            else:
                                return {"error": "No content returned from tool"}
                                
            except asyncio.TimeoutError:
                logger.warning(f"Tool call timeout (attempt {attempt + 1}/{max_retries})")
                if attempt < max_retries - 1:
                    await asyncio.sleep(retry_delay)
                else:
                    return {"error": "Tool call timed out after multiple attempts"}
            except Exception as e:
                logger.error(f"Tool call failed (attempt {attempt + 1}/{max_retries}): {e}")
                if attempt < max_retries - 1:
                    await asyncio.sleep(retry_delay)
                else:
                    return {"error": f"Tool call failed: {str(e)}"}
    
    async def close(self):
        """Clean up resources"""
        try:
            with self._lock:
                self.is_connected = False
                self.server_params = None
            print("✅ MCP resources cleaned up successfully")
        except Exception as e:
            print(f"❌ Error cleaning up MCP resources: {e}")

# Global MCP connection
mcp_connection = MCPConnection()

@tool
def create_metadata(text: str, file_name: str, label: str, classification_labels: dict = None):
    """
    Create a metadata dictionary for a given text document.
    
    Args:
        text: The document text content
        file_name: Name of the file
        label: Classification label for the document
        classification_labels: Optional dictionary mapping file names to classification labels
    
    Returns:
        dict: Metadata dictionary with document information
    """
    # Ưu tiên sử dụng nhãn từ classification_labels nếu có
    if classification_labels and isinstance(classification_labels, dict):
        print(f"\n🔎 CREATE_METADATA: Checking classification_labels for {file_name}")
        print(f"   - Classification labels: {classification_labels}")
        
        # Thử tìm theo tên file chính xác
        if file_name in classification_labels and classification_labels[file_name]:
            label = classification_labels[file_name]
            print(f"   - ✅ Using exact match from classification_labels: '{label}'")
        else:
            # Thử tìm theo tên file không phân biệt hoa thường
            file_name_lower = file_name.lower()
            for key, value in classification_labels.items():
                if isinstance(key, str) and key.lower() == file_name_lower and value:
                    label = value
                    print(f"   - ✅ Using case-insensitive match from classification_labels: '{label}'")
                    break
    
    # Clean up the label if it contains a file path
    if isinstance(label, str) and ':' in label:
        label = label.split(':')[-1].strip()
    
    # Ensure we have a valid label
    if not label or not isinstance(label, str) or label.lower() in ["không xác định", "chưa phân loại", "không có phân loại"]:
        # Xác định nhãn mặc định dựa trên tên file
        if "finance" in file_name.lower():
            label = "Tài chính"
        elif "report" in file_name.lower():
            label = "Báo cáo"
        else:
            label = "Tài liệu"
        print(f"   - ⚠️ Invalid label. Using default category based on filename: '{label}'")
    
    # Clean file name
    file_name = str(file_name).strip()
    if not file_name or file_name.lower() == 'unknown_file':
        file_name = f"document_{int(datetime.now().timestamp())}"
    
    # Create clean metadata structure
    metadata = {
        "file_name": file_name,
        "label": label,
        "content": text[:500] if text and len(text) > 500 else (text or ""),
        "total_characters": float(len(text)) if text else 0.0,
        "creation_date": _format_file_timestamp(
            timestamp=datetime.now().timestamp(), 
            include_time=True
        )
    }
    
    print(f"✅ Created metadata for {file_name} with label: {label}")
    return {"create_metadata_response": metadata}

@tool
def save_metadata_to_mcp(metadata: dict):
    """
    Save metadata to MCP server using proper MCP protocol.
    
    Args:
        metadata: Dictionary containing metadata to save
    
    Returns:
        dict: Result of the save operation
    """
    if not mcp_connection:
        return {"error": "MCP connection not initialized. Please initialize MCP first."}
    
    try:
        # Extract metadata from different possible structures
        if 'create_metadata_response' in metadata:
            # If it's a direct response from create_metadata
            create_resp = metadata['create_metadata_response']
            # Handle nested create_metadata_response if present
            if 'create_metadata_response' in create_resp:
                create_resp = create_resp['create_metadata_response']
                
            filename = create_resp.get('file_name', create_resp.get('filename', 'unknown_file'))
            label = create_resp.get('label', 'unclassified')
            content = create_resp.get('content', '')
            
            # Prepare clean additional metadata without duplicating fields
            additional_metadata = {k: v for k, v in create_resp.items() 
                                if k not in ['file_name', 'filename', 'label', 'content']}
            
            print(f"📝 Preparing to save metadata for {filename} with label {label}")
            print(f"📄 Content length: {len(content) if content else 0} chars")
        else:
            # If it's regular metadata
            filename = metadata.get('file_name', metadata.get('filename', 'unknown_file'))
            label = metadata.get('label', 'unclassified')
            content = metadata.get('content', '')
            additional_metadata = metadata
        
        # Debug output
        print(f"💾 Saving metadata for file: {filename}, label: {label}")
        
        # Call MCP server tool synchronously
        result = mcp_connection.call_tool_sync("save_metadata_to_json", {
            "filename": filename,
            "label": label,
            "content": content,
            "additional_metadata": additional_metadata
        })
        
        print(f"📊 Raw MCP result: {result}")
        
        # Check for explicit error field
        if isinstance(result, dict) and "error" in result:
            print(f"❌ MCP reported error: {result['error']}")
            return {
                "status": "error",
                "message": f"MCP server error: {result['error']}",
                "details": result.get("traceback", "")
            }
        
        # Handle string result (which is the expected format from MCP server)
        if isinstance(result, str):
            try:
                parsed_result = json.loads(result)
                print(f"✅ Successfully parsed JSON result")
                return {
                    "status": "success",
                    "message": "Metadata saved to MCP server successfully",
                    "metadata_id": parsed_result.get("metadata", {}).get("id", ""),
                    "result": parsed_result
                }
            except json.JSONDecodeError as e:
                print(f"⚠️ Failed to parse JSON: {e}")
        
        # Handle case where result is already a dict with result field
        if isinstance(result, dict) and "result" in result:
            try:
                if isinstance(result["result"], str):
                    parsed_inner = json.loads(result["result"])
                    print(f"✅ Successfully parsed nested JSON result")
                    return {
                        "status": "success",
                        "message": "Metadata saved to MCP server successfully",
                        "metadata_id": parsed_inner.get("metadata", {}).get("id", ""),
                        "result": parsed_inner
                    }
            except (json.JSONDecodeError, TypeError) as e:
                print(f"⚠️ Failed to parse nested JSON: {e}")
        
        # Return the raw result if all else fails
        print(f"ℹ️ Returning raw result")
        return {
            "status": "success",
            "message": "Metadata saved to MCP server successfully",
            "result": result
        }
        
    except Exception as e:
        import traceback
        return {
            "status": "error",
            "message": f"Failed to save metadata: {str(e)}",
            "details": traceback.format_exc()
        }

@tool
def search_metadata_in_mcp(filename: Optional[str] = None, label: Optional[str] = None):
    """
    Search for metadata in the MCP server.
    
    Args:
        filename: Optional filename to search for
        label: Optional label to search for
    
    Returns:
        dict: Search results from MCP server
    """
    try:
        # Debug output
        search_params = []
        if filename:
            search_params.append(f"filename='{filename}'")
        if label:
            search_params.append(f"label='{label}'")
        search_desc = " and ".join(search_params) if search_params else "all records"
        print(f"🔍 Searching metadata for {search_desc}")
        
        # Call MCP server tool synchronously
        result = mcp_connection.call_tool_sync("search_metadata", {
            "filename": filename,
            "label": label
        })
        
        print(f"📊 Raw search result: {result}")
        
        # Check for explicit error field
        if isinstance(result, dict) and "error" in result:
            print(f"❌ MCP search error: {result['error']}")
            return {
                "status": "error",
                "message": f"MCP server error: {result['error']}"
            }
        
        # Handle string result (which is the expected format from MCP server)
        if isinstance(result, str):
            try:
                parsed_result = json.loads(result)
                print(f"✅ Successfully parsed search JSON result")
                return parsed_result
            except json.JSONDecodeError as e:
                print(f"⚠️ Failed to parse search JSON: {e}")
        
        # Handle case where result is already a dict with result field
        if isinstance(result, dict) and "result" in result:
            try:
                if isinstance(result["result"], str):
                    parsed_inner = json.loads(result["result"])
                    print(f"✅ Successfully parsed nested search JSON result")
                    return parsed_inner
            except (json.JSONDecodeError, TypeError) as e:
                print(f"⚠️ Failed to parse nested search JSON: {e}")
        
        # Return the raw result if all else fails
        print(f"ℹ️ Returning raw search result")
        return result
        
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"❌ Search error: {e}\n{error_details}")
        return {
            "status": "error",
            "message": f"Failed to search metadata: {str(e)}",
            "details": error_details
        }

@tool
def get_metadata_from_mcp(metadata_id: str):
    """
    Get metadata from MCP server by ID.
    
    Args:
        metadata_id: ID of the metadata to retrieve
    
    Returns:
        dict: Metadata from MCP server
    """
    try:
        # Debug output
        print(f"🌐 Getting metadata with ID: {metadata_id}")
        
        # Call MCP server tool synchronously
        result = mcp_connection.call_tool_sync("get_metadata", {
            "metadata_id": metadata_id
        })
        
        print(f"📊 Raw get result: {result}")
        
        # Check for explicit error field
        if isinstance(result, dict) and "error" in result:
            print(f"❌ MCP get error: {result['error']}")
            return {
                "status": "error",
                "message": f"MCP server error: {result['error']}"
            }
        
        # Handle string result (which is the expected format from MCP server)
        if isinstance(result, str):
            try:
                parsed_result = json.loads(result)
                print(f"✅ Successfully parsed get JSON result")
                return parsed_result
            except json.JSONDecodeError as e:
                print(f"⚠️ Failed to parse get JSON: {e}")
        
        # Handle case where result is already a dict with result field
        if isinstance(result, dict) and "result" in result:
            try:
                if isinstance(result["result"], str):
                    parsed_inner = json.loads(result["result"])
                    print(f"✅ Successfully parsed nested get JSON result")
                    return parsed_inner
            except (json.JSONDecodeError, TypeError) as e:
                print(f"⚠️ Failed to parse nested get JSON: {e}")
        
        # Return the raw result if all else fails
        print(f"ℹ️ Returning raw get result")
        return result
        
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"❌ Get error: {e}\n{error_details}")
        return {
            "status": "error",
            "message": f"Failed to get metadata: {str(e)}",
            "details": error_details
        }

@tool
def list_all_metadata_from_mcp():
    """
    List all metadata from MCP server.
    
    Returns:
        dict: All metadata from MCP server
    """
    try:
        # Debug output
        print(f"📃 Listing all metadata from MCP server")
        
        # Call MCP server tool synchronously
        result = mcp_connection.call_tool_sync("list_all_metadata", {})
        
        print(f"📊 Raw list result: {result}")
        
        # Check for explicit error field
        if isinstance(result, dict) and "error" in result:
            print(f"❌ MCP list error: {result['error']}")
            return {
                "status": "error",
                "message": f"MCP server error: {result['error']}"
            }
        
        # Handle string result (which is the expected format from MCP server)
        if isinstance(result, str):
            try:
                parsed_result = json.loads(result)
                print(f"✅ Successfully parsed list JSON result")
                return parsed_result
            except json.JSONDecodeError as e:
                print(f"⚠️ Failed to parse list JSON: {e}")
        
        # Handle case where result is already a dict with result field
        if isinstance(result, dict) and "result" in result:
            try:
                if isinstance(result["result"], str):
                    parsed_inner = json.loads(result["result"])
                    print(f"✅ Successfully parsed nested list JSON result")
                    return parsed_inner
            except (json.JSONDecodeError, TypeError) as e:
                print(f"⚠️ Failed to parse nested list JSON: {e}")
        
        # Return the raw result if all else fails
        print(f"ℹ️ Returning raw list result")
        return result
        
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"❌ List error: {e}\n{error_details}")
        return {
            "status": "error",
            "message": f"Failed to list metadata: {str(e)}",
            "details": error_details
        }

class MetadataAgent(BaseAgent):
    """Metadata Agent backed by LangGraph with improved MCP integration."""
    
    @property
    def name(self):
        return "MetadataAgent"
    
    def __init__(self):
        super().__init__(
            agent_name='MetadataAgent',
            description='Tạo và quản lý metadata cho tài liệu',
            system_message=metadata_prompt,
            model=gemini,
            content_types=['text/plain', 'application/pdf', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/vnd.openxmlformats-officedocument.presentationml.presentation', '*/*']
        )
        self.mcp_initialized = False
        
        # Initialize graph for LangGraph integration
        self.graph = create_react_agent(
            self.model,
            checkpointer=memory,
            prompt=metadata_prompt,
            response_format=ResponseFormat,
            tools=[
                create_metadata, 
                save_metadata_to_mcp, 
                search_metadata_in_mcp, 
                get_metadata_from_mcp,
                list_all_metadata_from_mcp
            ],
            name="Metadata Agent",
        )
        
    # Synchronous MCP initialization
    def initialize_mcp_sync(self):
        """Synchronous MCP initialization"""
        if not self.mcp_initialized and MCP_AVAILABLE:
            print("🔌 Initializing MCP connection...")
            try:
                # Use the synchronous wrapper
                success = mcp_connection._run_async_in_thread(mcp_connection.start_server())
                if isinstance(success, dict) and "error" in success:
                    print(f"❌ MCP initialization failed: {success['error']}")
                    return False
                
                self.mcp_initialized = success
                if success:
                    print("✅ MCP connection initialized successfully")
                else:
                    print("❌ Failed to initialize MCP connection")
                return success
            except Exception as e:
                print(f"❌ Error initializing MCP: {e}")
                return False
        return self.mcp_initialized
        
    def extract_content_from_file(self, file_path):
        """Extract content directly from a file based on its extension"""
        try:
            if not os.path.exists(file_path):
                print(f"❌ File not found: {file_path}")
                return f"File not found: {file_path}"
                
            file_ext = os.path.splitext(file_path)[1].lower()
            print(f"📄 Extracting content from {file_path} (type: {file_ext})")
            
            # Extract based on file type
            if file_ext == '.pdf':
                try:
                    from agents.text_extraction_agent import extract_text_from_pdf
                    content = extract_text_from_pdf(file_path)
                except ImportError:
                    # Fallback implementation
                    import PyPDF2
                    with open(file_path, 'rb') as file:
                        reader = PyPDF2.PdfReader(file)
                        content = ""
                        for page in reader.pages:
                            content += page.extract_text() + "\n"
                print(f"✅ Extracted {len(content)} characters from PDF")
                
            elif file_ext in ['.docx', '.doc']:
                try:
                    from agents.text_extraction_agent import extract_text_from_word
                    content = extract_text_from_word(file_path)
                except ImportError:
                    # Fallback implementation
                    import docx
                    doc = docx.Document(file_path)
                    content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                print(f"✅ Extracted {len(content)} characters from Word document")
                
            elif file_ext in ['.pptx', '.ppt']:
                try:
                    from agents.text_extraction_agent import extract_text_from_powerpoint
                    content = extract_text_from_powerpoint(file_path)
                except ImportError:
                    # Fallback implementation
                    import pptx
                    presentation = pptx.Presentation(file_path)
                    content = ""
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                print(f"✅ Extracted {len(content)} characters from PowerPoint")
                
            else:
                # For other file types, try to read as text
                try:
                    with open(file_path, 'r', encoding='utf-8') as file:
                        content = file.read()
                    print(f"✅ Extracted {len(content)} characters from text file")
                except Exception as e:
                    print(f"❌ Cannot extract text from {file_path}: {e}")
                    content = f"Cannot extract text from this file type: {file_ext}"
            
            return content
        except Exception as e:
            print(f"❌ Error extracting content from {file_path}: {e}")
            return f"Error extracting content: {str(e)}"

    def invoke(self, query, metadata=None, sessionId=None):
        """
        Invoke the metadata agent to extract and save metadata.
        
        Args:
            query: The query to process
            metadata: Optional metadata dictionary
            sessionId: Session ID for the conversation
            
        Returns:
            str: Response from the agent
        """
        # Initialize MCP connection if needed
        global mcp_connection
        if not mcp_connection:
            try:
                from utils.mcp_connection import MCPConnection
                mcp_connection = MCPConnection()
                print("✅ MCP connection initialized")
            except Exception as e:
                print(f"❌ Failed to initialize MCP connection: {e}")
        
        # In ra toàn bộ metadata để debug
        print(f"\n🔍 FULL METADATA RECEIVED IN METADATA AGENT:")
        if metadata:
            print(f"Metadata type: {type(metadata)}")
            print(f"Metadata keys: {list(metadata.keys()) if isinstance(metadata, dict) else 'Not a dictionary'}")
            
            # In ra chi tiết về classification_labels nếu có
            if isinstance(metadata, dict) and 'classification_labels' in metadata:
                print(f"classification_labels type: {type(metadata['classification_labels'])}")
                print(f"classification_labels content: {metadata['classification_labels']}")
            else:
                print(f"⚠️ No classification_labels key in metadata")
                
            # In ra các thông tin quan trọng khác
            if isinstance(metadata, dict):
                for key in ['file_names', 'file_paths', 'is_multi_file', 'file_count']:
                    if key in metadata:
                        print(f"{key}: {metadata[key]}")
        else:
            print("⚠️ No metadata provided")
        
        # Debug log for metadata
        if metadata:
            print(f"\n📋 Received metadata in invoke:")
            print(f"- File name: {metadata.get('file_name', 'N/A')}")
            print(f"- Label: {metadata.get('label', 'N/A')}")
            content_len = len(metadata.get('content', '')) if metadata.get('content') else 0
            print(f"- Content length: {content_len} characters")
        else:
            print("⚠️ No metadata received in invoke method")
        
        # Prepare the input message
        input_data = {
            'messages': [('user', query)]
        }
    
        # Add metadata to the query if available
        if metadata:
            # Check if this is a multi-file metadata request
            is_multi_file = metadata.get('is_multi_file', False)
            file_paths = metadata.get('file_paths', [])
            
            if is_multi_file and file_paths and len(file_paths) > 1:
                # Handle multiple files
                try:
                    import os
                    file_names = metadata.get('file_names', [])
                    if not file_names and file_paths:
                        file_names = [os.path.basename(path) for path in file_paths]
                        
                    # Lấy classification_labels từ metadata để sử dụng cho từng file
                    classification_labels = metadata.get('classification_labels', {})
                    print(f"DEBUG: Available classification_labels for multi-file: {classification_labels}")
                    
                    # Lấy nhãn chung nếu có
                    general_label = metadata.get('label')
                    if general_label is None or general_label == 'None':
                        general_label = ""
                        print(f"⚠️ No general label provided for multi-file case, using empty label")
                    else:
                        print(f"✅ Using general label from metadata: '{general_label}'")
                        
                    # Ghi log số lượng nhãn phân loại có sẵn
                    if classification_labels:
                        print(f"✅ Found {len(classification_labels)} classification labels for {len(file_names)} files")
                    else:
                        print(f"⚠️ No classification labels found in metadata dictionary")
                    content = metadata.get('content', '')
                        
                    print(f"✅ Creating metadata for {len(file_paths)} files with general label: {general_label}")
                    print(f"Content length: {len(content)} characters")
                        
                    # Create metadata dictionary for the group
                    # Ensure we have at least some content
                    safe_content = content if content else f"Multiple files: {', '.join(file_names[:3])}{'...' if len(file_names) > 3 else ''} (no content extracted)"
                    metadata_dict = {
                        "file_name": metadata.get('file_name', 'multiple_files'),
                        "label": general_label,  # Sử dụng general_label thay vì label
                        "content": safe_content[:500] if len(safe_content) > 500 else safe_content,
                        "total_characters": float(len(safe_content)),
                        "creation_date": _format_file_timestamp(
                            timestamp=datetime.now().timestamp(), 
                            include_time=True
                        ),
                        "file_count": len(file_paths),
                        "file_names": file_names,
                        "classification_labels": classification_labels  # Thêm classification_labels vào metadata
                    }
                    
                    # Save metadata for each file
                    metadata_ids = []
                    for i, file_path in enumerate(file_paths):
                        file_name = file_names[i] if i < len(file_names) else os.path.basename(file_path)
                        
                        # Save metadata to MCP
                        if mcp_connection:
                            # First try to extract content directly from the file
                            try:
                                print(f"🔍 Attempting direct content extraction for {file_name}")
                                direct_content = self.extract_content_from_file(file_path)
                                if direct_content and len(direct_content) > 50:  # Ensure we got meaningful content
                                    file_content = direct_content
                                    print(f"📄 Successfully extracted content directly: {len(file_content)} characters")
                                else:
                                    # Fall back to passed content if direct extraction failed or returned minimal content
                                    print(f"⚠️ Direct extraction returned minimal content, checking alternatives")
                                    # Check if we have individual content for this file
                                    individual_contents = metadata.get('individual_contents', {})
                                    if file_name in individual_contents:
                                        file_content = individual_contents[file_name]
                                        print(f"📋 Using individual content from metadata: {len(file_content)} characters")
                                    else:
                                        file_content = content
                                        print(f"📋 Using shared content for file {file_name}")
                            except Exception as e:
                                print(f"❌ Error in direct extraction, falling back: {e}")
                                # Check if we have individual content for this file
                                individual_contents = metadata.get('individual_contents', {})
                                if file_name in individual_contents:
                                    file_content = individual_contents[file_name]
                                    print(f"📋 Using individual content from metadata: {len(file_content)} characters")
                                else:
                                    file_content = content
                                    print(f"📋 Using shared content for file {file_name}")
                                
                            # Ensure we have at least some content
                            safe_content = file_content if file_content else f"File: {file_name} (no content extracted)"
                            print(f"📋 Saving metadata for file {i+1}/{len(file_paths)}: {file_name}")
                            print(f"   - General Label: {general_label}")
                            print(f"   - Content length: {len(safe_content)} characters")
                            
                            # Parse the enhanced query to extract file-specific labels
                            # This is a more reliable approach than relying on metadata dictionaries
                            file_specific_label = None
                            
                            # Try to extract from the query which contains the file list with labels
                            if query and "DANH SÁCH FILES:" in query:
                                import re
                                # Look for different possible patterns in the enhanced query
                                # Pattern 1: "+ N. filename - Phân loại: label"
                                pattern1 = re.compile(fr"\+\s*\d+\.\s*{re.escape(file_name)}\s*-\s*Phân loại:\s*([^\n]+)")
                                # Pattern 2: "+ N. filename - label" (without "Phân loại:")
                                pattern2 = re.compile(fr"\+\s*\d+\.\s*{re.escape(file_name)}\s*-\s*([^\n]+)")
                                
                                # Try pattern 1 first (with "Phân loại:")
                                match = pattern1.search(query)
                                if match:
                                    file_specific_label = match.group(1).strip()
                                    print(f"   - Found label with pattern1 for {file_name}: {file_specific_label}")
                                else:
                                    # Try pattern 2 (without "Phân loại:")
                                    match = pattern2.search(query)
                                    if match:
                                        file_specific_label = match.group(1).strip()
                                        print(f"   - Found label with pattern2 for {file_name}: {file_specific_label}")
                                
                                # Add extra debug info
                                if not match:
                                    print(f"   - DEBUG: Could not find label pattern for {file_name} in query")
                                    # Print the relevant section of the query for debugging
                                    query_lines = query.split('\n')
                                    files_section_start = False
                                    for line in query_lines:
                                        if "DANH SÁCH FILES:" in line:
                                            files_section_start = True
                                            print(f"   - DEBUG: Found files section: {line}")
                                        elif files_section_start and file_name in line:
                                            print(f"   - DEBUG: Found file in query: {line}")
                                            break
                            
                            # If we found a specific label in the query, use it
                            if file_specific_label:
                                file_label = file_specific_label
                                print(f"   - Using label from query for {file_name}: {file_label}")
                            else:
                                # Fall back to checking in metadata dictionaries
                                file_label = general_label  # Default to the general label
                            
                            # Ưu tiên sử dụng classification_labels trước file_labels
                            classification_labels = metadata.get('classification_labels', {})
                            print(f"DEBUG: Available classification_labels: {classification_labels}")
                            print(f"DEBUG: Current file_name: {file_name}")
                            
                            # Try different ways to match the file name
                            found_label = False
                            
                            # 1. Direct match in classification_labels by file_name
                            if file_name in classification_labels:
                                file_label = classification_labels[file_name]
                                print(f"   - Using classification label for {file_name}: {file_label}")
                                found_label = True
                            # 2. Try with full file path
                            elif file_path in classification_labels:
                                file_label = classification_labels[file_path]
                                print(f"   - Using classification label for full path {file_path}: {file_label}")
                                found_label = True
                            # 3. Try case-insensitive match in classification_labels
                            else:
                                file_name_lower = file_name.lower()
                                for key, value in classification_labels.items():
                                    if isinstance(key, str):
                                        # Check if key is filename or path
                                        if key.lower() == file_name_lower or (os.path.basename(key).lower() == file_name_lower):
                                            file_label = value
                                            print(f"   - Using case-insensitive classification label match for {file_name} -> {key}: {file_label}")
                                            found_label = True
                                            break
                            
                            # If still not found, fall back to checking in file_labels
                            if not found_label:
                                file_labels = metadata.get('file_labels', {})
                                print(f"DEBUG: Available file_labels: {file_labels}")
                                
                                # Direct match
                                if file_name in file_labels:
                                    file_label = file_labels[file_name]
                                    print(f"   - Using specific label for {file_name}: {file_label}")
                                    found_label = True
                                else:
                                    # Try case-insensitive match
                                    for key, value in file_labels.items():
                                        if key.lower() == file_name_lower:
                                            file_label = value
                                            print(f"   - Using case-insensitive match for {file_name} -> {key}: {file_label}")
                                            found_label = True
                                            break
                            
                            if not found_label:
                                print(f"DEBUG: No specific label found for {file_name} in either dictionary, using default: {file_label}")
                            
                            metadata_id = str(uuid.uuid4())
                            # Calculate actual content length
                            actual_content_length = len(safe_content)
                            print(f"   - Actual content length: {actual_content_length} characters")
                            
                            # Prepare additional metadata specific to this file
                            additional_meta = {
                                "total_characters": actual_content_length,  # Use actual integer value, not float
                                "creation_date": _format_file_timestamp(
                                    timestamp=datetime.now().timestamp(), 
                                    include_time=True
                                ),
                                "file_path": file_path,
                            }
                            
                            # Only add group information if this is part of a multi-file group
                            if len(file_paths) > 1:
                                additional_meta["is_part_of_group"] = True
                                additional_meta["group_size"] = len(file_paths)
                                # Only include file_names of other files in the group, not this file
                                other_files = [f for f in file_names if f != file_name]
                                if other_files:
                                    additional_meta["related_files"] = other_files
                            
                            # Kiểm tra lần cuối và làm sạch nhãn phân loại trước khi lưu
                            print(f"   - 🔍 CHECKING LABEL: '{file_label}'")
                            
                            # Thử tìm lại trong classification_labels một lần nữa để đảm bảo chúng ta có nhãn chính xác nhất
                            classification_labels = metadata.get('classification_labels', {})
                            original_label = file_label
                            
                            # In ra chi tiết về classification_labels
                            print(f"\n🔎 DETAILED CLASSIFICATION LABELS DEBUG:")
                            print(f"   - Classification labels type: {type(classification_labels)}")
                            print(f"   - Classification labels content: {classification_labels}")
                            print(f"   - File name being processed: '{file_name}'")
                            print(f"   - Is file_name in classification_labels? {file_name in classification_labels}")
                            if file_name in classification_labels:
                                print(f"   - Value for '{file_name}': '{classification_labels[file_name]}'")
                                print(f"   - Value type: {type(classification_labels[file_name])}")
                                print(f"   - Is value truthy? {bool(classification_labels[file_name])}")
                            
                            # Kiểm tra lại trong classification_labels
                            if classification_labels:
                                if file_name in classification_labels and classification_labels[file_name]:
                                    file_label = classification_labels[file_name]
                                    print(f"   - ✅ PRIORITY: Using exact classification label from metadata: '{file_label}'")
                                else:
                                    print(f"   - ⚠️ No exact match found or value is empty. Trying case-insensitive match.")
                                    # Thử tìm theo tên file không phân biệt hoa thường
                                    file_name_lower = file_name.lower()
                                    for key, value in classification_labels.items():
                                        print(f"   - Checking key: '{key}', value: '{value}', match? {isinstance(key, str) and key.lower() == file_name_lower and value}")
                                        if isinstance(key, str) and key.lower() == file_name_lower and value:
                                            file_label = value
                                            print(f"   - ✅ PRIORITY: Using case-insensitive classification label: '{file_label}'")
                                            break
                            
                            # Chỉ kiểm tra và sử dụng nhãn mặc định nếu không tìm thấy nhãn hợp lệ trong classification_labels
                            # Kiểm tra nếu nhãn trống hoặc None
                            if not file_label or file_label == "None" or file_label.strip() == "":
                                # Xác định nhãn mặc định dựa trên tên file
                                if "finance" in file_name.lower():
                                    file_label = "Tài chính"
                                elif "report" in file_name.lower():
                                    file_label = "Báo cáo"
                                else:
                                    file_label = "Tài liệu"
                                print(f"   - ⚠️ Empty label detected. Using default category based on filename: '{file_label}'")
                            # Kiểm tra nếu nhãn có dạng đường dẫn file
                            elif file_label.startswith('C:\\') or file_label.startswith('/') or 'ĐƯỜNG DẪN' in file_label or file_label.startswith('path:'):
                                print(f"   - ⚠️ FINAL CHECK: Label looks like a file path: '{file_label}'. Using default category.")
                                # Xác định nhãn mặc định dựa trên tên file
                                if "finance" in file_name.lower():
                                    file_label = "Tài chính"
                                elif "report" in file_name.lower():
                                    file_label = "Báo cáo"
                                else:
                                    file_label = "Tài liệu"
                                print(f"   - 🏷️ Using default category: '{file_label}'")
                            # Kiểm tra nếu nhãn quá dài (có thể là nội dung file)
                            elif len(file_label) > 50:
                                print(f"   - ⚠️ Label too long ({len(file_label)} chars), likely content not label. Using default category.")
                                # Xác định nhãn mặc định dựa trên tên file
                                if "finance" in file_name.lower():
                                    file_label = "Tài chính"
                                elif "report" in file_name.lower():
                                    file_label = "Báo cáo"
                                else:
                                    file_label = "Tài liệu"
                                print(f"   - 🏷️ Using default category: '{file_label}'")
                                
                            # Ghi log nếu nhãn đã thay đổi
                            if original_label != file_label:
                                print(f"   - 🔄 Label changed from '{original_label}' to '{file_label}'")
                            else:
                                print(f"   - ✅ Label unchanged: '{file_label}'")
                            
                                
                            # In ra nhãn cuối cùng sẽ được lưu
                            print(f"   - 📌 FINAL LABEL for {file_name}: '{file_label}'")
                            

                            # Tạo metadata trực tiếp thay vì gọi hàm create_metadata
                            print(f"\n🔎 CREATING METADATA WITH CLASSIFICATION_LABELS")
                            
                            # Kiểm tra lại classification_labels một lần nữa
                            if classification_labels and isinstance(classification_labels, dict):
                                print(f"   - Classification labels available: {len(classification_labels)} entries")
                                if file_name in classification_labels:
                                    print(f"   - Found exact match for {file_name} in classification_labels: {classification_labels[file_name]}")
                            else:
                                print(f"   - No valid classification_labels dictionary available")
                            
                            # Tạo metadata trực tiếp
                            metadata_obj = {
                                "id": str(uuid.uuid4()),
                                "filename": file_name,
                                "file_name": file_name,  # Đảm bảo có cả hai trường
                                "label": file_label,
                                "content": safe_content[:500] if safe_content and len(safe_content) > 500 else (safe_content or ""),
                                "total_characters": float(len(safe_content)) if safe_content else 0.0,
                                "created_at": datetime.now().isoformat(),
                                "updated_at": datetime.now().isoformat(),
                                "creation_date": _format_file_timestamp(
                                    timestamp=datetime.now().timestamp(), 
                                    include_time=True
                                ),
                                "additional_metadata": additional_meta
                            }
                            
                            # In ra metadata trước khi lưu
                            print(f"   - 💾 Saving metadata with label: '{metadata_obj.get('label')}' (type: {type(metadata_obj.get('label'))})")
                            print(f"   - 💾 Metadata details: {metadata_obj}")
                            
                            # Save metadata to MCP
                            result = mcp_connection.call_tool_sync("save_metadata_to_json", metadata_obj)
                            
                            # Extract metadata ID and handle errors
                            metadata_id = None
                            if isinstance(result, dict):
                                # Check for explicit error
                                if 'error' in result:
                                    print(f"❌ MCP error for file {file_name}: {result['error']}")
                                # Try to extract ID from different possible structures
                                elif 'metadata' in result and 'id' in result['metadata']:
                                    metadata_id = result['metadata']['id']
                                elif 'id' in result:
                                    metadata_id = result['id']
                                
                                # Debug print the result structure
                                print(f"📊 MCP result keys: {list(result.keys())}")
                            
                            if metadata_id:
                                metadata_ids.append(metadata_id)
                                print(f"✅ Metadata saved for file {file_name} with ID: {metadata_id}")
                    
                    if metadata_ids:
                        # Return the first ID as the primary ID
                        primary_id = metadata_ids[0]
                        return f"✅ Đã tạo và lưu metadata thành công cho {len(metadata_ids)} files. Metadata ID chính: {primary_id}"
                    else:
                        print("❌ No metadata IDs were returned from MCP. This might indicate a server issue.")
                        return f"❌ Không thể lưu metadata cho các files. Đã thử lưu {len(file_paths)} files với nhãn '{label}'."
                except Exception as e:
                    import traceback
                    print(f"❌ Error creating metadata for multiple files: {e}")
                    print(traceback.format_exc())
                    return f"❌ Lỗi khi tạo metadata cho nhiều files: {str(e)}"
            
            # Handle single file
            elif metadata.get('content'):
                try:
                    # Get metadata parameters from the metadata dict
                    file_paths = metadata.get('file_paths', [])
                    file_names = metadata.get('file_names', [])
                    label = metadata.get('label')
                    file_name = metadata.get('file_name', 'unknown_file')
                    file_label = ""
                    
                    # Thử lấy nhãn từ nhiều nguồn khác nhau
                    print(f"\n🔍 SINGLE FILE LABEL EXTRACTION for {file_name}:")
                    
                    # Ưu tiên lấy từ classification_labels trước tiên
                    classification_labels = metadata.get('classification_labels', {})
                    print(f"   - Available classification_labels: {classification_labels}")
                    original_label = file_label
                    
                    if classification_labels:
                        # Thử tìm theo tên file chính xác
                        if file_name in classification_labels and classification_labels[file_name]:
                            file_label = classification_labels[file_name]
                            print(f"   - ✅ PRIORITY: Using exact match in classification_labels: '{file_label}'")
                        else:
                            # Thử tìm theo tên file không phân biệt hoa thường
                            file_name_lower = file_name.lower()
                            for key, value in classification_labels.items():
                                if isinstance(key, str) and key.lower() == file_name_lower and value:
                                    file_label = value
                                    print(f"   - ✅ PRIORITY: Using case-insensitive match in classification_labels: '{file_label}'")
                                    break
                    
                    # Nếu không tìm thấy trong classification_labels, thử lấy từ label chung
                    if (not file_label or file_label == "None" or file_label.strip() == "") and label and label != 'None':
                        file_label = label
                        print(f"   - Using general label as fallback: '{file_label}'")
                        
                    # Ghi log nếu nhãn đã thay đổi
                    if original_label != file_label:
                        print(f"   - 🔄 Label changed from '{original_label}' to '{file_label}'")
                    else:
                        print(f"   - ✅ Label unchanged: '{file_label}'")
                    
                    # 3. Kiểm tra và xử lý nhãn
                    if not file_label or file_label == "None" or file_label.strip() == "":
                        # Xác định nhãn mặc định dựa trên tên file
                        if "finance" in file_name.lower():
                            file_label = "Tài chính"
                        elif "report" in file_name.lower():
                            file_label = "Báo cáo"
                        else:
                            file_label = "Tài liệu"
                        print(f"   - ⚠️ No valid label found. Using default category based on filename: '{file_label}'")
                    # Kiểm tra nếu nhãn có dạng đường dẫn file
                    elif file_label.startswith('C:\\') or file_label.startswith('/') or 'ĐƯỜNG DẪN' in file_label or file_label.startswith('path:'):
                        print(f"   - ⚠️ Label looks like a file path: '{file_label}'. Using default category.")
                        # Xác định nhãn mặc định dựa trên tên file
                        if "finance" in file_name.lower():
                            file_label = "Tài chính"
                        elif "report" in file_name.lower():
                            file_label = "Báo cáo"
                        else:
                            file_label = "Tài liệu"
                        print(f"   - 🏷️ Using default category: '{file_label}'")
                    
                    print(f"   - 📌 FINAL LABEL for single file {file_name}: '{file_label}'")
                    
                    content = metadata.get('content', '')
                    
                    file_name = metadata.get('file_name', 'unknown_file')
                    print(f"✅ Directly creating metadata for {file_name} with label: {file_label}")
                    print(f"Content length: {len(content)} characters")
                    
                 
                    # Ensure we have at least some content
                    safe_content = content if content else f"File: {file_name} (no content extracted)"
                    found_label = False
                    # file_label đã được định nghĩa ở trên
                    if metadata.get('classification_labels'):
                        classification_labels = metadata.get('classification_labels', {})
                        if file_name in classification_labels:
                            file_label = classification_labels[file_name]
                            found_label = True
                    # In ra nhãn cuối cùng sẽ được lưu
                    print(f"   - 📌 FINAL LABEL for single file {file_name}: '{file_label}'")
                    
                    # Sử dụng hàm create_metadata để tạo metadata với classification_labels
                    print(f"\n🔎 USING CREATE_METADATA WITH CLASSIFICATION_LABELS FOR SINGLE FILE")
                    metadata_result = create_metadata(
                        text=safe_content,
                        file_name=file_name,
                        label=file_label,
                        classification_labels=classification_labels
                    )
                    
                    # Lấy metadata từ kết quả của create_metadata
                    if isinstance(metadata_result, dict) and 'create_metadata_response' in metadata_result:
                        metadata_obj = metadata_result['create_metadata_response']
                        # Thêm các trường bổ sung
                        metadata_obj["id"] = str(uuid.uuid4())
                        metadata_obj["filename"] = file_name  # Đảm bảo filename được đặt đúng
                        metadata_obj["file_path"] = file_paths[0] if file_paths else ""
                        metadata_obj["classification_source"] = "classification_labels" if found_label else "query_extraction"
                    else:
                        # Fallback nếu create_metadata không trả về kết quả mong đợi
                        metadata_obj = {
                            "id": str(uuid.uuid4()),
                            "filename": file_name,
                            "label": file_label,
                            "content": safe_content[:500] if len(safe_content) > 500 else safe_content,
                            "total_characters": float(len(safe_content)),
                            "creation_date": _format_file_timestamp(
                                timestamp=datetime.now().timestamp(), 
                                include_time=True
                            ),
                            "file_path": file_paths[0] if file_paths else "",
                            "classification_source": "classification_labels" if found_label else "query_extraction"
                        }
                        print(f"   - ⚠️ Fallback to direct metadata creation due to unexpected create_metadata result")
                    
                    print(f"   - Final label for {file_name}: '{metadata_obj.get('label')}' (type: {type(metadata_obj.get('label'))})")
                    print(f"   - Label source: {metadata_obj.get('classification_source', 'unknown')}")
                    
                    # Save metadata to MCP
                    if mcp_connection:
                        # In ra metadata trước khi lưu
                        print(f"📋 Saving metadata for single file: {file_name}")
                        print(f"   - Label: {metadata_obj.get('label')}")
                        print(f"   - Content length: {len(metadata_obj.get('content', ''))} characters")
                        print(f"   - 💾 Metadata details: {metadata_obj}")
                        
                        # Save metadata to MCP
                        result = mcp_connection.call_tool_sync("save_metadata_to_json", metadata_obj)
                        
                        # Extract metadata ID and handle errors
                        metadata_id = None
                        if isinstance(result, dict):
                            # Check for explicit error
                            if 'error' in result:
                                print(f"❌ MCP error for file {file_name}: {result['error']}")
                            # Try to extract ID from different possible structures
                            elif 'metadata' in result and 'id' in result['metadata']:
                                metadata_id = result['metadata']['id']
                            elif 'id' in result:
                                metadata_id = result['id']
                            
                            # Debug print the result structure
                            print(f"📊 MCP result keys: {list(result.keys())}")
                        
                        if metadata_id:
                            print(f"✅ Metadata saved with ID: {metadata_id}")
                            return f"✅ Đã tạo và lưu metadata thành công cho file {file_name}. Metadata ID: {metadata_id}"
                        else:
                            print(f"⚠️ Failed to save metadata for file {file_name}, no ID returned")
                            return f"❌ Không thể lưu metadata cho file {file_name}. Vui lòng kiểm tra lại kết nối MCP server."
                except Exception as e:
                    import traceback
                    print(f"❌ Error creating metadata directly: {e}")
                    print(traceback.format_exc())
        
        # If direct creation failed or wasn't attempted, continue with normal flow
        config = {
            'configurable': {
                'thread_id': sessionId,
                'metadata': metadata  # Include in config for tool access
            },
            'recursion_limit': 50
        }
        
        try:
            # Pass both input_data and config to the graph
            response = self.graph.invoke(input_data, config)
            
            # Extract the last message from the response
            if isinstance(response, dict) and 'messages' in response:
                last_message = response['messages'][-1]
            elif isinstance(response, (list, tuple)) and len(response) > 0:
                last_message = response[-1]
            else:
                last_message = response
            
            # Extract content from the message
            if hasattr(last_message, 'content'):
                content = last_message.content
            else:
                content = str(last_message)
                
            return content
        except Exception as e:
            import traceback
            error_msg = f"Error processing request: {str(e)}"
            print(f"❌ {error_msg}")
            print(traceback.format_exc())
            return error_msg
    
    async def stream(self, query, sessionId, task_id) -> AsyncIterable[Dict[str, Any]]:
        """Stream responses from the agent."""
        # Initialize MCP if not done
        if not self.mcp_initialized:
            success = await mcp_connection.start_server()
            self.mcp_initialized = success
            
        inputs = {'messages': [('user', query)]}
        config = {'configurable': {'thread_id': sessionId}, 'recursion_limit': 50}

        try:
            for item in self.graph.stream(inputs, config, stream_mode='values'):
                message = item['messages'][-1]
                if isinstance(message, AIMessage):
                    yield {
                        'response_type': 'text',
                        'is_task_complete': False,
                        'require_user_input': False,
                        'content': message.content,
                    }
            
            yield {
                'response_type': 'data', 
                'is_task_complete': True,
                'require_user_input': False,
                'content': message.content if isinstance(message, AIMessage) else str(message)
            }
        except Exception as e:
            yield {
                'response_type': 'error',
                'is_task_complete': True,
                'require_user_input': False,
                'content': f"Error in streaming: {str(e)}"
            }
                    
    async def ainvoke(self, input_data, config=None):
        """Required method for langgraph-supervisor compatibility"""
        # Initialize MCP if not done
        if not self.mcp_initialized:
            await mcp_connection.start_server()
            self.mcp_initialized = True
            
        if config is None:
            config = {}
            
        try:
            messages = input_data.get("messages", [])
            if messages and len(messages) > 0:
                last_message = messages[-1]
                
                if hasattr(last_message, 'content'):
                    query = last_message.content
                elif isinstance(last_message, dict) and 'content' in last_message:
                    query = last_message['content']
                elif isinstance(last_message, tuple) and len(last_message) == 2:
                    query = last_message[1]
                else:
                    query = str(last_message)
                    
                session_id = config.get("configurable", {}).get("thread_id", "default")
                result = self.invoke(query, session_id)
                
                if isinstance(result, dict):
                    result = str(result)
                    
                from langchain_core.messages import AIMessage
                return {"messages": [AIMessage(content=result)]}
        except Exception as e:
            from langchain_core.messages import AIMessage
            return {"messages": [AIMessage(content=f"Error in ainvoke: {str(e)}")]}

    def __del__(self):
        """Cleanup MCP connection on deletion"""
        try:
            # Don't create new event loop in destructor
            pass
        except Exception as e:
            print(f"Warning: Error during cleanup: {e}")


if __name__ == "__main__":
    def test_metadata_agent():
        print("🧪 Testing MetadataAgent with improved MCP server integration...")
        
        # Create agent
        metadata_agent = MetadataAgent()
        
        # Initialize MCP connection
        print("🔌 Initializing MCP connection...")
        success = metadata_agent.initialize_mcp_sync()
        
        if not success:
            print("❌ Failed to connect to MCP server. Please check:")
            print("1. MCP server file exists")
            print("2. MCP server dependencies are installed")
            print("3. No firewall/permission issues")
            return
        
        # Test prompt
        prompt = """
        Hãy tạo metadata cho tài liệu sau:
        - File Name: Chain_of_thought.pdf
        - Label: Chain of Thought
        - Text: This is a document about the chain of thought. The document is a PDF file containing research about reasoning methods in artificial intelligence. It discusses various approaches to breaking down complex problems into step-by-step reasoning processes.
        
        Sau đó, lưu metadata vào MCP server và báo cáo kết quả chi tiết.
        """
        
        try:
            print("📝 Calling MetadataAgent to process request...")
            result = metadata_agent.invoke(query=prompt, sessionId="test_session")
            print("\n✅ Agent result:")
            print(result)
            
            
        except Exception as e:
            import traceback
            print(f"❌ Error during testing: {e}")
            print(traceback.format_exc())
        finally:
            # Cleanup
            print("\n🧹 Cleaning up...")
            # Note: cleanup is handled automatically
    
    # Run the synchronous test
    test_metadata_agent()