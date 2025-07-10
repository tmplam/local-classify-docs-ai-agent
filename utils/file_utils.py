import os
import sys
import re
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import logging

logger = logging.getLogger("agent_logger")

def read_file_content(file_path, max_chars=50000):
    """
    Đọc nội dung của file dựa trên định dạng của nó.
    Hỗ trợ các định dạng: .txt, .pdf, .docx, .pptx
    
    Args:
        file_path (str): Đường dẫn đến file cần đọc
        max_chars (int): Số ký tự tối đa trả về
        
    Returns:
        str: Nội dung của file hoặc chuỗi rỗng nếu không đọc được
    """
    if not os.path.exists(file_path):
        logger.error(f"File không tồn tại: {file_path}")
        return ""
    
    try:
        file_ext = os.path.splitext(file_path)[1].lower()
        
        # Xử lý file text
        if file_ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(max_chars)
                
        # Xử lý file PDF
        elif file_ext == ".pdf":
            text = ""
            try:
                doc = fitz.open(file_path)
                for page in doc:
                    text += page.get_text()
                    if len(text) > max_chars:
                        break
                return text[:max_chars]
            except Exception as e:
                logger.error(f"Lỗi khi đọc file PDF {file_path}: {str(e)}")
                return ""
                
        # Xử lý file Word
        elif file_ext == ".docx":
            try:
                doc = Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
                return text[:max_chars]
            except Exception as e:
                logger.error(f"Lỗi khi đọc file Word {file_path}: {str(e)}")
                return ""
                
        # Xử lý file PowerPoint
        elif file_ext == ".pptx":
            try:
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text[:max_chars]
            except Exception as e:
                logger.error(f"Lỗi khi đọc file PowerPoint {file_path}: {str(e)}")
                return ""
                
        # Các định dạng khác không được hỗ trợ
        else:
            logger.warning(f"Định dạng file không được hỗ trợ: {file_ext}")
            return ""
            
    except Exception as e:
        logger.error(f"Lỗi khi đọc file {file_path}: {str(e)}")
        return ""

def clean_invalid_unicode(text):
    """
    Xóa các ký tự Unicode không hợp lệ từ văn bản
    
    Args:
        text (str): Văn bản cần làm sạch
        
    Returns:
        str: Văn bản đã được làm sạch
    """
    if not text:
        return ""
        
    # Loại bỏ các ký tự điều khiển ngoại trừ tab, newline và carriage return
    text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
    
    # Thay thế các ký tự không hợp lệ bằng khoảng trắng
    text = ''.join(ch if ord(ch) < 0x10000 else ' ' for ch in text)
    
    return text
